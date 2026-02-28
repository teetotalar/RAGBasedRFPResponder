# src/ingest_kb.py

import os
import uuid
import fitz  # PDF
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct


COLLECTION_NAME = "knowledge_base"
VECTOR_SIZE = 384  # all-MiniLM-L6-v2
CHUNK_WORD_SIZE = 300  # better retrieval granularity


# Load embedding model once
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# Local Qdrant storage
qdrant = QdrantClient(path="qdrant_storage")


# ----------------------------------------
# TEXT EXTRACTION FUNCTIONS
# ----------------------------------------

def extract_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def extract_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_html(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f, "html.parser")
        return soup.get_text(separator="\n")


def extract_markdown(file_path):
    return extract_txt(file_path)


# ----------------------------------------
# UNIVERSAL EXTRACTOR
# ----------------------------------------

def extract_text(file_path):

    ext = file_path.lower().split(".")[-1]

    if ext == "pdf":
        return extract_pdf(file_path)

    elif ext == "docx":
        return extract_docx(file_path)

    elif ext == "pptx":
        return extract_pptx(file_path)

    elif ext == "txt":
        return extract_txt(file_path)

    elif ext in ["html", "htm"]:
        return extract_html(file_path)

    elif ext == "md":
        return extract_markdown(file_path)

    else:
        print(f"⚠ Unsupported format skipped: {file_path}")
        return ""


# ----------------------------------------
# CHUNKING
# ----------------------------------------

def chunk_text(text, chunk_size=CHUNK_WORD_SIZE):

    words = text.split()
    chunks = []

    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)

    return chunks


# ----------------------------------------
# COLLECTION MANAGEMENT
# ----------------------------------------

def create_collection_if_not_exists():

    collections = qdrant.get_collections().collections
    existing = [c.name for c in collections]

    if COLLECTION_NAME not in existing:
        qdrant.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(
                size=VECTOR_SIZE,
                distance=Distance.COSINE
            )
        )
        print(f"Created collection: {COLLECTION_NAME}")


# ----------------------------------------
# MAIN INGESTION FUNCTION
# ----------------------------------------

def ingest_folder(folder_path):

    if not os.path.exists(folder_path):
        print(f"Folder not found: {folder_path}")
        return

    create_collection_if_not_exists()

    for file in os.listdir(folder_path):

        file_path = os.path.join(folder_path, file)

        if not os.path.isfile(file_path):
            continue

        print(f"Ingesting: {file}")

        text = extract_text(file_path)

        if not text.strip():
            print(f"⚠ Skipped empty file: {file}")
            continue

        chunks = chunk_text(text)

        points = []

        for chunk in chunks:
            vector = embedding_model.encode(chunk).tolist()

            points.append(
                PointStruct(
                    id=str(uuid.uuid4()),
                    vector=vector,
                    payload={
                        "text": chunk,
                        "source_file": file
                    }
                )
            )

        qdrant.upsert(
            collection_name=COLLECTION_NAME,
            points=points
        )

    print("Ingestion complete.")